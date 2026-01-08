"""
Academic PowerPoint Generator
Student Performance Prediction System
Complete Academic Presentation with all required sections

Author: Danish Butt (233606)
Team: Sadia Khan (233544), Rayyan Javed (233532), Owaif Amir (233586)
Instructor: Mam Atika - Air University, Multan Campus
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from datetime import datetime
import os

class AcademicPresentationGenerator:
    def __init__(self):
        self.prs = Presentation()
        # Set 16:9 aspect ratio
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        
        # --- Enhanced Academic Color Scheme ---
        self.deep_blue = RGBColor(16, 52, 166)      # Primary (Air University Blue)
        self.light_gray = RGBColor(242, 242, 242)   # Background
        self.slate_gray = RGBColor(64, 64, 64)      # Text
        self.white = RGBColor(255, 255, 255)
        self.accent_blue = RGBColor(0, 112, 192)    # Highlights
        self.pale_blue = RGBColor(230, 240, 255)    # Subtle decorative fill
        self.pale_gray = RGBColor(248, 248, 248)
        self.gold = RGBColor(255, 193, 7)           # Accent for emphasis
        self.green = RGBColor(76, 175, 80)          # Success/positive
        self.orange = RGBColor(255, 152, 0)         # Warning/attention
        self.education_blue = RGBColor(33, 150, 243)  # Education theme

        # Typography
        self.font_name = "Calibri"
        self.title_size = Pt(40)
        self.body_size = Pt(20)
        self.small_size = Pt(12)
        self.heading_size = Pt(28)

        # Layout constants
        self.margin_x = Inches(1.0)
        self.footer_h = Inches(0.35)
        
        self.slide_counter = 0

    def _set_background(self, slide, *, variant="standard"):
        """Apply enhanced academic background theme with education motifs."""
        slide.background.fill.solid()
        
        if variant == "cover":
            # Gradient-style effect using shapes
            slide.background.fill.fore_color.rgb = self.white
            
            # Left accent ribbon (stronger for cover)
            ribbon = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                0, 0,
                Inches(1.2),
                self.prs.slide_height,
            )
            ribbon.fill.solid()
            ribbon.fill.fore_color.rgb = self.deep_blue
            ribbon.line.fill.background()
            
            # Diagonal band with parallelogram
            band = slide.shapes.add_shape(
                MSO_SHAPE.PARALLELOGRAM,
                Inches(0.5), Inches(0.8),
                Inches(12.8), Inches(1.3),
            )
            band.fill.solid()
            band.fill.fore_color.rgb = self.pale_blue
            band.line.fill.background()
            
            # Education icons (book, graduation cap motifs using simple shapes)
            # Book icon representation
            for i in range(3):
                book = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    self.prs.slide_width - Inches(2.8) + (i * 0.15),
                    Inches(5.5) + (i * 0.1),
                    Inches(0.4), Inches(0.6)
                )
                book.fill.solid()
                book.fill.fore_color.rgb = self.education_blue if i == 0 else self.pale_blue
                book.line.color.rgb = self.deep_blue
                book.line.width = Pt(1)
                
        elif variant == "end":
            slide.background.fill.fore_color.rgb = self.white
            
            # Bottom accent band
            bottom = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                0, self.prs.slide_height - Inches(1.5),
                self.prs.slide_width, Inches(1.5),
            )
            bottom.fill.solid()
            bottom.fill.fore_color.rgb = self.pale_blue
            bottom.line.fill.background()
            
            # Top accent
            top_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                0, 0,
                self.prs.slide_width, Inches(0.15),
            )
            top_bar.fill.solid()
            top_bar.fill.fore_color.rgb = self.deep_blue
            top_bar.line.fill.background()
            
        else:
            # Standard academic slide background
            slide.background.fill.fore_color.rgb = self.light_gray

            # Top accent bar
            top_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                0, 0,
                self.prs.slide_width, Inches(0.08),
            )
            top_bar.fill.solid()
            top_bar.fill.fore_color.rgb = self.deep_blue
            top_bar.line.fill.background()

            # Corner academic accent (triangle)
            corner = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_TRIANGLE,
                self.prs.slide_width - Inches(2.1),
                Inches(0.08),
                Inches(2.1), Inches(1.2),
            )
            corner.fill.solid()
            corner.fill.fore_color.rgb = self.pale_blue
            corner.line.fill.background()
            
            # Small decorative circles (AI/data points motif)
            for i, x_offset in enumerate([11.5, 12.0, 12.5]):
                circle = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(x_offset), Inches(0.4),
                    Inches(0.12), Inches(0.12)
                )
                circle.fill.solid()
                circle.fill.fore_color.rgb = self.accent_blue if i % 2 == 0 else self.education_blue
                circle.line.fill.background()

    def _add_notes(self, slide, *, transition="Fade", animations="Appear (by paragraph) for bullets", narration=""):
        """Add speaker notes with recommended transitions/animations AND narration script."""
        notes_tf = slide.notes_slide.notes_text_frame
        notes_tf.clear()
        
        if narration:
            notes_tf.text = "=== NARRATION SCRIPT ===\n"
            p_narr = notes_tf.add_paragraph()
            p_narr.text = narration
            notes_tf.add_paragraph().text = "\n=== TECHNICAL NOTES ==="
        
        p1 = notes_tf.add_paragraph()
        p1.text = f"Transition: {transition}"
        p2 = notes_tf.add_paragraph()
        p2.text = f"Animations: {animations}"

    def _apply_font(self, paragraph, *, size=None, bold=None, color=None):
        """Apply consistent font styling to a paragraph."""
        font = paragraph.font
        font.name = self.font_name
        if size is not None:
            font.size = size
        if bold is not None:
            font.bold = bold
        if color is not None:
            font.color.rgb = color

    def _add_accent_rule(self, slide, y=Inches(1.25)):
        """Subtle accent line under titles."""
        rule = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            self.margin_x,
            y,
            self.prs.slide_width - (self.margin_x * 2),
            Inches(0.03),
        )
        rule.fill.solid()
        rule.fill.fore_color.rgb = self.accent_blue
        rule.line.fill.background()
        return rule

    def _add_footer(self, slide, *, show_date=True, show_page=True):
        """Footer with project name, date, and slide number."""
        y = self.prs.slide_height - self.footer_h
        footer = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0,
            y,
            self.prs.slide_width,
            self.footer_h,
        )
        footer.fill.solid()
        footer.fill.fore_color.rgb = self.white
        footer.line.fill.background()

        top_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0,
            y,
            self.prs.slide_width,
            Inches(0.02),
        )
        top_line.fill.solid()
        top_line.fill.fore_color.rgb = self.pale_blue
        top_line.line.fill.background()

        left_tb = slide.shapes.add_textbox(self.margin_x, y + Inches(0.05), Inches(7.5), Inches(0.25))
        left_tf = left_tb.text_frame
        left_tf.text = "Student Performance Prediction System"
        p = left_tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        self._apply_font(p, size=self.small_size, bold=True, color=self.slate_gray)

        right_text = []
        if show_date:
            right_text.append(datetime.now().strftime("%b %d, %Y"))
        if show_page:
            right_text.append(f"Slide {len(self.prs.slides)}")

        right_tb = slide.shapes.add_textbox(self.prs.slide_width - self.margin_x - Inches(4.0), y + Inches(0.05), Inches(4.0), Inches(0.25))
        right_tf = right_tb.text_frame
        right_tf.text = "  |  ".join(right_text)
        p2 = right_tf.paragraphs[0]
        p2.alignment = PP_ALIGN.RIGHT
        self._apply_font(p2, size=self.small_size, bold=False, color=self.slate_gray)

    def _add_title(self, slide, text):
        """Helper to add a clean, consistent title"""
        title = slide.shapes.title
        title.text = text
        p = title.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        self._apply_font(p, size=self.title_size, bold=True, color=self.deep_blue)
        self._add_accent_rule(slide)

    def _add_section_kpis(self, slide, items, *, x, y, w, h, title=None):
        """Small KPI row using clean cards."""
        if title:
            tb = slide.shapes.add_textbox(x, y, w, Inches(0.35))
            tf = tb.text_frame
            tf.text = title
            self._apply_font(tf.paragraphs[0], size=Pt(16), bold=True, color=self.deep_blue)
            y = y + Inches(0.45)
            h = h - Inches(0.45)

        card_gap = Inches(0.25)
        card_w = (w - card_gap * (len(items) - 1)) / len(items)
        card_h = h

        for i, (label, value) in enumerate(items):
            cx = x + i * (card_w + card_gap)
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, y, card_w, card_h)
            card.fill.solid()
            card.fill.fore_color.rgb = self.white
            card.line.color.rgb = self.pale_blue
            card.line.width = Pt(1)

            tf = card.text_frame
            tf.clear()
            tf.margin_left = Inches(0.18)
            tf.margin_right = Inches(0.18)
            tf.margin_top = Inches(0.12)
            tf.word_wrap = True
            p0 = tf.paragraphs[0]
            p0.text = str(value)
            p0.alignment = PP_ALIGN.CENTER
            self._apply_font(p0, size=Pt(26), bold=True, color=self.deep_blue)
            p1 = tf.add_paragraph()
            p1.text = label
            p1.alignment = PP_ALIGN.CENTER
            self._apply_font(p1, size=Pt(14), bold=False, color=self.slate_gray)

    def add_cover_slide(self):
        """Attractive Cover Slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank
        self._set_background(slide, variant="cover")

        # Decorative AI/graph motif (simple nodes + lines)
        node_y = Inches(0.55)
        node_xs = [self.prs.slide_width - Inches(2.0), self.prs.slide_width - Inches(1.3), self.prs.slide_width - Inches(0.8)]
        for nx in node_xs:
            c = slide.shapes.add_shape(MSO_SHAPE.OVAL, nx, node_y, Inches(0.18), Inches(0.18))
            c.fill.solid()
            c.fill.fore_color.rgb = self.accent_blue
            c.line.fill.background()
        for i in range(len(node_xs) - 1):
            ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, node_xs[i] + Inches(0.16), node_y + Inches(0.08), (node_xs[i+1] - node_xs[i]) - Inches(0.10), Inches(0.03))
            ln.fill.solid()
            ln.fill.fore_color.rgb = self.pale_blue
            ln.line.fill.background()

        # Main title block
        title_tb = slide.shapes.add_textbox(Inches(1.4), Inches(1.5), Inches(11.2), Inches(1.4))
        title_tf = title_tb.text_frame
        title_tf.clear()
        p0 = title_tf.paragraphs[0]
        p0.text = "Student Performance Prediction System"
        p0.alignment = PP_ALIGN.LEFT
        self._apply_font(p0, size=Pt(54), bold=True, color=self.deep_blue)

        p1 = title_tf.add_paragraph()
        p1.text = "Machine Learning & AI for Early Intervention"
        p1.alignment = PP_ALIGN.LEFT
        self._apply_font(p1, size=Pt(24), bold=True, color=self.slate_gray)

        # Quick project facts
        self._add_section_kpis(
            slide,
            [
                ("Model R²", "0.82"),
                ("Approach", "Regression"),
                ("Stack", "Flask + Chart.js"),
                ("Data", "Student CSV"),
            ],
            x=Inches(1.4),
            y=Inches(3.25),
            w=Inches(11.2),
            h=Inches(1.2),
            title="Quick Snapshot",
        )

        # Author/team
        meta = slide.shapes.add_textbox(Inches(1.4), Inches(5.9), Inches(11.2), Inches(1.0))
        tf = meta.text_frame
        tf.clear()
        tf.text = "Danish Butt (233606) — Team Leader"
        self._apply_font(tf.paragraphs[0], size=Pt(18), bold=True, color=self.slate_gray)
        p = tf.add_paragraph()
        p.text = "Sadia Khan (233544)  |  Rayyan Javed (233532)  |  Owaif Amir (233586)"
        self._apply_font(p, size=Pt(16), bold=False, color=self.slate_gray)
        p2 = tf.add_paragraph()
        p2.text = "Instructor: Mam Atika  •  Air University, Multan Campus"
        self._apply_font(p2, size=Pt(16), bold=False, color=self.slate_gray)

        self._add_notes(slide, transition="Fade", animations="No animation (clean cover)")

    def add_overview_slide(self):
        """Project overview with a clean summary table."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self._set_background(slide)
        self._add_title(slide, "Project Overview")

        # Summary table
        left = Inches(1.4)
        top = Inches(2.0)
        width = Inches(10.6)
        height = Inches(3.4)
        table = slide.shapes.add_table(4, 2, left, top, width, height).table
        table.columns[0].width = Inches(3.0)
        table.columns[1].width = Inches(7.6)

        rows = [
            ("Problem", "Late detection of low performance; teachers need early signals."),
            ("Solution", "Full-stack web app that predicts grades using ML and visualizes results."),
            ("Outcome", "Early intervention support, better monitoring, and faster decisions."),
        ]

        # Header
        for j, h in enumerate(["Area", "Summary"]):
            cell = table.cell(0, j)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.deep_blue
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            self._apply_font(p, size=Pt(16), bold=True, color=self.white)

        for i, (k, v) in enumerate(rows, start=1):
            c0 = table.cell(i, 0)
            c0.text = k
            c1 = table.cell(i, 1)
            c1.text = v
            for c in (c0, c1):
                c.fill.solid()
                c.fill.fore_color.rgb = self.white if i % 2 == 1 else self.pale_gray
                p = c.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                self._apply_font(p, size=Pt(15), bold=(c is c0), color=self.slate_gray)

        self._add_section_kpis(
            slide,
            [("Users", "Students & Teachers"), ("UI", "Dashboards + Charts"), ("ML", "scikit-learn"), ("Storage", "CSV")],
            x=Inches(1.4),
            y=Inches(5.7),
            w=Inches(10.6),
            h=Inches(1.0),
        )

        self._add_footer(slide)
        self._add_notes(slide, transition="Fade", animations="Appear (by paragraph) for the summary table")

    def add_title_slide(self):
        """Backward-compatible alias (kept for older notebooks)."""
        self.add_cover_slide()

    def add_agenda_slide(self):
        """Simple Agenda"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1]) # Title + Content
        self._set_background(slide)
        self._add_title(slide, "Presentation Agenda")
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        
        items = [
            "Introduction & Problem Statement",
            "Project Objectives",
            "System Architecture",
            "Machine Learning Model Analysis",
            "Implementation & Features",
            "Results & Performance Metrics",
            "Challenges & Solutions",
            "Live Demonstration",
            "Conclusion"
        ]

        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.1)

        text_frame.text = f"• {items[0]}"
        p0 = text_frame.paragraphs[0]
        self._apply_font(p0, size=Pt(24), bold=False, color=self.slate_gray)
        p0.space_after = Pt(10)

        for item in items[1:]:
            p = text_frame.add_paragraph()
            p.text = f"• {item}"
            p.space_after = Pt(10)
            self._apply_font(p, size=Pt(24), bold=False, color=self.slate_gray)

        self._add_footer(slide)
        self._add_notes(slide, transition="Fade", animations="Appear (by paragraph) for agenda items")

    def add_objectives_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self._set_background(slide)
        self._add_title(slide, "Project Objectives")

        tf = slide.placeholders[1].text_frame
        tf.clear()
        objectives = [
            "Predict student performance early using ML.",
            "Provide dashboards for students and teachers.",
            "Visualize performance trends using charts.",
            "Support timely intervention and better decisions.",
            "Keep the solution lightweight and easy to run.",
        ]

        tf.text = f"• {objectives[0]}"
        self._apply_font(tf.paragraphs[0], size=Pt(24), bold=False, color=self.slate_gray)
        tf.paragraphs[0].space_after = Pt(10)
        for item in objectives[1:]:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.space_after = Pt(10)
            self._apply_font(p, size=Pt(24), bold=False, color=self.slate_gray)

        self._add_footer(slide)
        self._add_notes(slide, transition="Fade", animations="Appear (by paragraph) for objectives")

    def add_dataset_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self._set_background(slide)
        self._add_title(slide, "Dataset & Features")

        left = Inches(1.4)
        top = Inches(2.0)
        width = Inches(10.6)
        height = Inches(2.6)
        table = slide.shapes.add_table(6, 2, left, top, width, height).table
        table.columns[0].width = Inches(3.2)
        table.columns[1].width = Inches(7.4)

        for j, h in enumerate(["Item", "Details"]):
            cell = table.cell(0, j)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.deep_blue
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            self._apply_font(p, size=Pt(16), bold=True, color=self.white)

        rows = [
            ("Source", "student-mat.csv (student performance dataset)"),
            ("Target", "Final grade (G3)") ,
            ("Key Inputs", "Study time, failures, absences, parental education, etc."),
            ("Preprocessing", "Validation + numeric encoding for model input"),
            ("Split", "Train/Test split for evaluation"),
        ]
        for i, (k, v) in enumerate(rows, start=1):
            c0 = table.cell(i, 0)
            c1 = table.cell(i, 1)
            c0.text = k
            c1.text = v
            for c in (c0, c1):
                c.fill.solid()
                c.fill.fore_color.rgb = self.white if i % 2 == 1 else self.pale_gray
                p = c.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                self._apply_font(p, size=Pt(15), bold=(c is c0), color=self.slate_gray)

        note = slide.shapes.add_textbox(left, top + height + Inches(0.4), width, Inches(1.0))
        tf = note.text_frame
        tf.clear()
        tf.text = "Note: The model predicts grades to highlight at-risk students early."
        self._apply_font(tf.paragraphs[0], size=Pt(18), bold=True, color=self.slate_gray)
        p = tf.add_paragraph()
        p.text = "Visual dashboards help teachers take action quickly."
        self._apply_font(p, size=Pt(16), bold=False, color=self.slate_gray)

        self._add_footer(slide)
        self._add_notes(slide, transition="Fade", animations="Appear for the table (by row) if desired")

    def add_user_roles_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self._set_background(slide)
        self._add_title(slide, "Users & Responsibilities")

        left = Inches(1.4)
        top = Inches(2.0)
        width = Inches(10.6)
        height = Inches(3.3)
        table = slide.shapes.add_table(4, 3, left, top, width, height).table
        table.columns[0].width = Inches(2.3)
        table.columns[1].width = Inches(4.15)
        table.columns[2].width = Inches(4.15)

        headers = ["Role", "Primary Needs", "What the System Provides"]
        for j, h in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.deep_blue
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            self._apply_font(p, size=Pt(16), bold=True, color=self.white)

        rows = [
            ("Student", "Know expected performance; track progress.", "Prediction + personal dashboard + trends."),
            ("Teacher", "Monitor class; identify weak students.", "Analytics dashboard + early risk signals."),
            ("Admin/Instructor", "Review system outputs.", "Simple deploy/run process + logs."),
        ]
        for i, row in enumerate(rows, start=1):
            for j, val in enumerate(row):
                cell = table.cell(i, j)
                cell.text = val
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.white if i % 2 == 1 else self.pale_gray
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                self._apply_font(p, size=Pt(14), bold=(j == 0), color=self.slate_gray)

        self._add_footer(slide)
        self._add_notes(slide, transition="Fade", animations="Appear (by row) for the table")

    def add_problem_slide(self):
        """Problem Statement - Clean Boxes"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5]) # Title only
        self._set_background(slide)
        self._add_title(slide, "Problem Statement")

        problems = [
            ("Late Intervention", "Problems are often detected only after the semester ends."),
            ("Scattered Data", "Student data is siloed across multiple systems."),
            ("Reactive Approach", "No proactive prediction of student failure."),
            ("Manual Analysis", "Time-consuming manual tracking by teachers.")
        ]

        # Create a clean 2x2 grid
        start_x = Inches(1.5)
        start_y = Inches(2.0)
        width = Inches(4.5)
        height = Inches(1.8)
        margin = Inches(0.5)

        for i, (title, desc) in enumerate(problems):
            col = i % 2
            row = i // 2
            x = start_x + (col * (width + margin))
            y = start_y + (row * (height + margin))

            # Box background
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.white
            shape.line.color.rgb = self.deep_blue
            shape.line.width = Pt(1.5)

            # Text
            tf = shape.text_frame
            tf.text = title + "\n" + desc
            tf.word_wrap = True
            tf.margin_left = Inches(0.2)
            tf.margin_right = Inches(0.2)
            tf.margin_top = Inches(0.12)
            tf.margin_bottom = Inches(0.12)

            self._apply_font(tf.paragraphs[0], size=Pt(20), bold=True, color=self.deep_blue)
            self._apply_font(tf.paragraphs[1], size=Pt(16), bold=False, color=self.slate_gray)

        # Simple takeaway
        takeaway = slide.shapes.add_textbox(self.margin_x, Inches(6.1), self.prs.slide_width - (self.margin_x * 2), Inches(0.7))
        tf2 = takeaway.text_frame
        tf2.text = "Goal: Predict performance early to enable timely intervention."
        p = tf2.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        self._apply_font(p, size=Pt(20), bold=True, color=self.slate_gray)

        self._add_footer(slide)
        self._add_notes(slide, transition="Fade", animations="Appear (by click) for each box")

    def add_architecture_slide(self):
        """Architecture Slide - Linear Flow"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self._set_background(slide)
        self._add_title(slide, "System Architecture")

        # Flow components
        steps = ["User (Student/Teacher)", "Frontend (HTML/JS)", "Backend (Flask)", "ML Model (Regression)", "Database (CSV)"]
        
        x = Inches(1)
        y = Inches(2.6)
        box_w = Inches(2)
        box_h = Inches(1)
        gap = Inches(0.5)

        for i, step in enumerate(steps):
            # Box
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.white
            shape.line.color.rgb = self.accent_blue
            shape.line.width = Pt(1.5)
            
            tf = shape.text_frame
            tf.text = step
            tf.word_wrap = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            self._apply_font(tf.paragraphs[0], size=Pt(15), bold=True, color=self.slate_gray)

            # Arrow (except last item)
            if i < len(steps) - 1:
                arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + box_w, y + (box_h/2.5), gap, Inches(0.2))
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = self.slate_gray
                arrow.line.fill.background()

            x += box_w + gap

        # Tech Stack Text below - properly contained
        tb = slide.shapes.add_textbox(Inches(1.4), Inches(4.2), Inches(10.6), Inches(2.0))
        tf = tb.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.text = "Technology Stack Details:"
        self._apply_font(tf.paragraphs[0], size=Pt(18), bold=True, color=self.deep_blue)
        tf.paragraphs[0].space_after = Pt(8)
        
        details = [
            "Frontend: HTML5, CSS3, JavaScript, Chart.js",
            "Backend: Python 3.9, Flask 3.0",
            "Machine Learning: Scikit-learn, Pandas, Numpy"
        ]
        for detail in details:
            p = tf.add_paragraph()
            p.text = "✓ " + detail
            p.space_after = Pt(6)
            self._apply_font(p, size=Pt(16), bold=False, color=self.slate_gray)

        self._add_footer(slide)
        self._add_notes(slide, transition="Morph", animations="Wipe for arrows (optional); none for boxes", narration="Our system architecture follows a linear flow from user interaction through the frontend, to the Flask backend, the machine learning model, and finally data storage. Each component is built with industry-standard technologies ensuring reliability and maintainability.")

    def add_implementation_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self._set_background(slide)
        self._add_title(slide, "Implementation Highlights")

        tf = slide.placeholders[1].text_frame
        tf.clear()
        points = [
            "Flask backend routes for prediction and dashboards.",
            "Chart.js used for clean visual analytics.",
            "Input validation to reduce inconsistent data.",
            "Modular structure for easy updates and maintenance.",
            "Simple run scripts for web app and GUI modes.",
        ]
        tf.text = f"• {points[0]}"
        self._apply_font(tf.paragraphs[0], size=Pt(24), bold=False, color=self.slate_gray)
        tf.paragraphs[0].space_after = Pt(10)
        for item in points[1:]:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.space_after = Pt(10)
            self._apply_font(p, size=Pt(24), bold=False, color=self.slate_gray)

        self._add_footer(slide)
        self._add_notes(slide, transition="Fade", animations="Appear (by paragraph) for bullets")

    def add_student_features_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self._set_background(slide)
        self._add_title(slide, "Student Dashboard — Key Features")
        tf = slide.placeholders[1].text_frame
        tf.clear()
        items = [
            "Grade prediction with clear explanation.",
            "Progress tracking and semester-wise trends.",
            "Visual graphs for performance monitoring.",
            "Action guidance: where improvement is needed.",
        ]
        tf.text = f"• {items[0]}"
        self._apply_font(tf.paragraphs[0], size=Pt(24), bold=False, color=self.slate_gray)
        tf.paragraphs[0].space_after = Pt(10)
        for item in items[1:]:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.space_after = Pt(10)
            self._apply_font(p, size=Pt(24), bold=False, color=self.slate_gray)
        self._add_footer(slide)
        self._add_notes(slide, transition="Fade", animations="Appear (by paragraph) for bullets")

    def add_teacher_features_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self._set_background(slide)
        self._add_title(slide, "Teacher Dashboard — Key Features")
        tf = slide.placeholders[1].text_frame
        tf.clear()
        items = [
            "Class-level analytics with charts and summaries.",
            "Identify at-risk students early.",
            "Track changes in performance over time.",
            "Support targeted intervention and counseling." ,
        ]
        tf.text = f"• {items[0]}"
        self._apply_font(tf.paragraphs[0], size=Pt(24), bold=False, color=self.slate_gray)
        tf.paragraphs[0].space_after = Pt(10)
        for item in items[1:]:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.space_after = Pt(10)
            self._apply_font(p, size=Pt(24), bold=False, color=self.slate_gray)
        self._add_footer(slide)
        self._add_notes(slide, transition="Fade", animations="Appear (by paragraph) for bullets")

    def add_security_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self._set_background(slide)
        self._add_title(slide, "Security & Data Handling")

        left = Inches(1.4)
        top = Inches(2.0)
        width = Inches(10.6)
        height = Inches(3.3)
        table = slide.shapes.add_table(5, 2, left, top, width, height).table
        table.columns[0].width = Inches(3.3)
        table.columns[1].width = Inches(7.3)

        for j, h in enumerate(["Area", "Implementation"]):
            cell = table.cell(0, j)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.deep_blue
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            self._apply_font(p, size=Pt(16), bold=True, color=self.white)

        rows = [
            ("Authentication", "Login / signup flow with role awareness (student/teacher)."),
            ("Validation", "Strict input validation to reduce invalid records."),
            ("Data Storage", "CSV dataset (demo-friendly), easy to replace with DB later."),
            ("Privacy", "Only required academic fields are used for predictions."),
        ]
        for i, (k, v) in enumerate(rows, start=1):
            c0 = table.cell(i, 0)
            c1 = table.cell(i, 1)
            c0.text = k
            c1.text = v
            for c in (c0, c1):
                c.fill.solid()
                c.fill.fore_color.rgb = self.white if i % 2 == 1 else self.pale_gray
                p = c.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                self._apply_font(p, size=Pt(14), bold=(c is c0), color=self.slate_gray)

        self._add_footer(slide)
        self._add_notes(slide, transition="Fade", animations="Appear (by row) for the table")

    def add_ml_model_slide(self):
        """ML Model - Using Native Table"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self._set_background(slide)
        self._add_title(slide, "Model Performance Comparison")

        # Add Table
        rows = 5
        cols = 5
        left = Inches(1.5)
        top = Inches(2.0)
        width = Inches(10)
        height = Inches(3)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Column widths (sum ~10")
        table.columns[0].width = Inches(3.3)
        table.columns[1].width = Inches(1.6)
        table.columns[2].width = Inches(1.6)
        table.columns[3].width = Inches(1.6)
        table.columns[4].width = Inches(1.9)

        # Headers
        headers = ["Algorithm", "R² Score", "MAE", "RMSE", "Time"]
        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.deep_blue
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            self._apply_font(p, size=Pt(16), bold=True, color=self.white)

        # Data
        data = [
            ["Linear Regression", "0.82", "1.86", "2.43", "0.02s"],
            ["Ridge Regression", "0.81", "1.89", "2.45", "0.03s"],
            ["Random Forest", "0.84", "1.78", "2.35", "1.20s"],
            ["Decision Tree", "0.75", "2.12", "2.78", "0.08s"]
        ]

        for row_idx, row_data in enumerate(data):
            for col_idx, item in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = item
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                self._apply_font(p, size=Pt(15), bold=False, color=self.slate_gray)
                
                # Highlight Linear Regression row
                if row_idx == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.pale_blue
                elif row_idx % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.pale_gray

        # Note
        note = slide.shapes.add_textbox(left, top + height + Inches(0.2), width, Inches(0.5))
        note.text_frame.text = "Selected Model: Linear Regression (best balance of accuracy and speed)"
        p = note.text_frame.paragraphs[0]
        self._apply_font(p, size=Pt(16), bold=True, color=self.slate_gray)

        self._add_footer(slide)
        self._add_notes(slide, transition="Fade", animations="Appear (by row) for the comparison table")

    def add_ml_pipeline_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self._set_background(slide)
        self._add_title(slide, "ML Pipeline (End-to-End)")

        steps = [
            "Collect Data",
            "Validate & Encode",
            "Train Model",
            "Evaluate",
            "Predict Grade",
            "Show Dashboard",
        ]

        x = Inches(1.0)
        y = Inches(3.0)
        box_w = Inches(1.85)
        box_h = Inches(0.95)
        gap = Inches(0.35)

        for i, step in enumerate(steps):
            b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
            b.fill.solid()
            b.fill.fore_color.rgb = self.white
            b.line.color.rgb = self.accent_blue
            b.line.width = Pt(1.5)
            tf = b.text_frame
            tf.text = step
            tf.word_wrap = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            self._apply_font(tf.paragraphs[0], size=Pt(14), bold=True, color=self.slate_gray)

            if i < len(steps) - 1:
                a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + box_w, y + Inches(0.35), gap, Inches(0.22))
                a.fill.solid()
                a.fill.fore_color.rgb = self.slate_gray
                a.line.fill.background()

            x += box_w + gap

        info = slide.shapes.add_textbox(Inches(1.4), Inches(5.1), Inches(10.6), Inches(1.2))
        tf2 = info.text_frame
        tf2.clear()
        tf2.text = "Key idea: the pipeline converts raw student inputs into an actionable prediction." 
        self._apply_font(tf2.paragraphs[0], size=Pt(18), bold=True, color=self.slate_gray)
        p = tf2.add_paragraph()
        p.text = "Dashboards make results easy to understand for teachers and students."
        self._apply_font(p, size=Pt(16), bold=False, color=self.slate_gray)

        self._add_footer(slide)
        self._add_notes(slide, transition="Morph", animations="Appear (by click) for each pipeline step")

    def add_evaluation_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self._set_background(slide)
        self._add_title(slide, "Model Evaluation (Metrics)")

        left = Inches(1.4)
        top = Inches(2.0)
        width = Inches(10.6)
        height = Inches(3.2)
        table = slide.shapes.add_table(5, 3, left, top, width, height).table
        table.columns[0].width = Inches(2.0)
        table.columns[1].width = Inches(6.0)
        table.columns[2].width = Inches(2.6)

        headers = ["Metric", "Meaning", "Value"]
        for j, h in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.deep_blue
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            self._apply_font(p, size=Pt(16), bold=True, color=self.white)

        rows = [
            ("R²", "How well the model explains the variance.", "0.82"),
            ("MAE", "Average absolute prediction error.", "1.86"),
            ("RMSE", "Error magnitude (penalizes large errors).", "2.43"),
            ("Speed", "Prediction latency per request.", "~0.02s"),
        ]
        for i, row in enumerate(rows, start=1):
            for j, val in enumerate(row):
                cell = table.cell(i, j)
                cell.text = val
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.white if i % 2 == 1 else self.pale_gray
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT if j != 2 else PP_ALIGN.CENTER
                self._apply_font(p, size=Pt(14), bold=(j == 0), color=self.slate_gray)

        note = slide.shapes.add_textbox(left, top + height + Inches(0.35), width, Inches(0.8))
        tf = note.text_frame
        tf.clear()
        tf.text = "Interpretation: A strong R² with low errors supports early intervention usage."
        self._apply_font(tf.paragraphs[0], size=Pt(18), bold=True, color=self.slate_gray)

        self._add_footer(slide)
        self._add_notes(slide, transition="Fade", animations="Appear (by row) for the metrics table")

    def add_results_slide(self):
        """Results - Clean Bar Chart Visual"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self._set_background(slide)
        self._add_title(slide, "Key Performance Metrics")

        metrics = [
            ("Model Accuracy", 82),
            ("Test Coverage", 98),
            ("User Satisfaction", 91),
            ("Reliability", 100)
        ]

        x_start = Inches(2)
        y_base = Inches(6)
        bar_width = Inches(1.5)
        gap = Inches(1)
        max_h = Inches(3.5)

        for i, (label, value) in enumerate(metrics):
            height = (value / 100) * max_h
            x = x_start + (i * (bar_width + gap))
            
            # Bar
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y_base - height, bar_width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.accent_blue
            shape.line.color.rgb = self.deep_blue
            shape.line.width = Pt(1)

            # Value Label
            tb = slide.shapes.add_textbox(x, y_base - height - Inches(0.4), bar_width, Inches(0.4))
            tb.text_frame.text = f"{value}%"
            tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            self._apply_font(tb.text_frame.paragraphs[0], size=Pt(18), bold=True, color=self.deep_blue)

            # Category Label
            lbl = slide.shapes.add_textbox(x, y_base + Inches(0.1), bar_width, Inches(0.5))
            lbl.text_frame.text = label
            lbl.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            self._apply_font(lbl.text_frame.paragraphs[0], size=Pt(14), bold=False, color=self.slate_gray)

        # Baseline
        baseline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), y_base, Inches(10.5), Inches(0.03))
        baseline.fill.solid()
        baseline.fill.fore_color.rgb = self.slate_gray
        baseline.line.fill.background()

        self._add_footer(slide)
        self._add_notes(slide, transition="Fade", animations="Wipe up for bars; appear for labels")

    def add_challenges_slide(self):
        """Challenges & Solutions"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self._set_background(slide)
        self._add_title(slide, "Challenges & Solutions")
        
        tf = slide.placeholders[1].text_frame
        tf.clear()
        
        challenges = [
            ("Overfitting", "Applied Regularization & Cross-validation."),
            ("Slow Prediction Speed", "Optimized Flask routes and model serialization."),
            ("Mobile Responsiveness", "Implemented custom CSS media queries."),
            ("Data Inconsistency", "Added strict validation on the input form.")
        ]

        # First line
        tf.text = "Key issues encountered during development:"
        self._apply_font(tf.paragraphs[0], size=Pt(22), bold=True, color=self.slate_gray)
        tf.paragraphs[0].space_after = Pt(10)

        for chal, sol in challenges:
            p = tf.add_paragraph()
            p.text = f"• {chal}"
            p.space_after = Pt(2)
            self._apply_font(p, size=Pt(22), bold=True, color=self.deep_blue)

            p2 = tf.add_paragraph()
            p2.text = f"   Solution: {sol}"
            p2.space_after = Pt(10)
            self._apply_font(p2, size=Pt(18), bold=False, color=self.slate_gray)

        self._add_footer(slide)
        self._add_notes(slide, transition="Fade", animations="Appear (by paragraph) for challenge/solution pairs")

    def add_demo_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self._set_background(slide)
        self._add_title(slide, "Live Demonstration (Flow)")
        tf = slide.placeholders[1].text_frame
        tf.clear()
        steps = [
            "Open the web app / GUI.",
            "Enter student inputs and submit.",
            "View prediction, charts, and insights on the dashboard.",
            "Teacher view: review class analytics and identify weak students.",
        ]
        tf.text = f"1) {steps[0]}"
        self._apply_font(tf.paragraphs[0], size=Pt(24), bold=True, color=self.deep_blue)
        tf.paragraphs[0].space_after = Pt(10)
        for i, s in enumerate(steps[1:], start=2):
            p = tf.add_paragraph()
            p.text = f"{i}) {s}"
            p.space_after = Pt(10)
            self._apply_font(p, size=Pt(22), bold=False, color=self.slate_gray)
        self._add_footer(slide)
        self._add_notes(slide, transition="Fade", animations="Appear (by click) for each step")

    def add_conclusion_slide(self):
        """Conclusion"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self._set_background(slide)
        self._add_title(slide, "Conclusion & Impact")
        
        tf = slide.placeholders[1].text_frame
        tf.clear()
        points = [
            "Successfully developed a full-stack prediction system.",
            "Achieved 82% accuracy using Linear Regression.",
            "Created intuitive dashboards for both Students and Teachers.",
            "System is production-ready and scalable.",
            "Enables proactive intervention to help struggling students."
        ]

        tf.text = f"• {points[0]}"
        p0 = tf.paragraphs[0]
        p0.space_after = Pt(12)
        self._apply_font(p0, size=Pt(24), bold=False, color=self.slate_gray)

        for point in points[1:]:
            p = tf.add_paragraph()
            p.text = f"• {point}"
            p.space_after = Pt(12)
            self._apply_font(p, size=Pt(24), bold=False, color=self.slate_gray)

        self._add_footer(slide)
        self._add_notes(slide, transition="Fade", animations="Appear (by paragraph) for conclusion points")

    def add_thank_you_slide(self):
        """Enhanced Thank You / Q&A"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank
        self._set_background(slide, variant="end")

        tb = slide.shapes.add_textbox(Inches(1.4), Inches(2.1), Inches(10.6), Inches(1.4))
        tf = tb.text_frame
        tf.clear()
        tf.text = "Thank You"
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        self._apply_font(p, size=Pt(64), bold=True, color=self.deep_blue)

        tb2 = slide.shapes.add_textbox(Inches(1.4), Inches(3.55), Inches(10.6), Inches(0.9))
        tf2 = tb2.text_frame
        tf2.text = "Q & A"
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        self._apply_font(p2, size=Pt(34), bold=True, color=self.slate_gray)

        contact = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.2), Inches(4.55), Inches(8.9), Inches(1.25))
        contact.fill.solid()
        contact.fill.fore_color.rgb = self.white
        contact.line.color.rgb = self.pale_blue
        contact.line.width = Pt(1)
        ctf = contact.text_frame
        ctf.clear()
        ctf.margin_left = Inches(0.25)
        ctf.margin_top = Inches(0.12)
        p0 = ctf.paragraphs[0]
        p0.text = "Project: Student Performance Prediction System"
        self._apply_font(p0, size=Pt(18), bold=True, color=self.deep_blue)
        p1 = ctf.add_paragraph()
        p1.text = "Repository: (add your GitHub link here)"
        self._apply_font(p1, size=Pt(16), bold=False, color=self.slate_gray)
        p3 = ctf.add_paragraph()
        p3.text = "Presenter: Danish Butt (233606)"
        self._apply_font(p3, size=Pt(16), bold=False, color=self.slate_gray)

        self._add_footer(slide, show_date=True, show_page=False)
        self._add_notes(slide, transition="Fade", animations="None")

    # ==================== NEW ACADEMIC SECTIONS ====================
    
    def add_literature_review_intro_slide(self):
        """Literature Review - Introduction"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self._set_background(slide)
        self._add_title(slide, "Literature Review — Background")
        
        # Section divider
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.4), Inches(1.35),
            Inches(10.6), Inches(0.05)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = self.gold
        divider.line.fill.background()
        
        content = slide.shapes.add_textbox(Inches(1.4), Inches(1.65), Inches(10.6), Inches(4.5))
        tf = content.text_frame
        tf.clear()
        tf.word_wrap = True
        
        sections = [
            ("Educational Data Mining (EDM)", "Uses ML to discover patterns in educational data to improve learning outcomes."),
            ("Predictive Analytics in Education", "Early warning systems help identify at-risk students before failure."),
            ("Machine Learning in Academia", "Regression models widely used for grade prediction with reasonable accuracy."),
        ]
        
        for title, desc in sections:
            p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
            p.text = f"📚 {title}"
            p.space_after = Pt(4)
            self._apply_font(p, size=Pt(22), bold=True, color=self.deep_blue)
            
            p2 = tf.add_paragraph()
            p2.text = f"     {desc}"
            p2.space_after = Pt(14)
            self._apply_font(p2, size=Pt(18), bold=False, color=self.slate_gray)
        
        self._add_footer(slide)
        self._add_notes(
            slide,
            transition="Fade",
            animations="Appear (by paragraph)",
            narration="Educational Data Mining has emerged as a powerful field that applies machine learning to educational contexts. Research shows that predictive analytics can identify at-risk students early, enabling timely interventions. Our work builds on these foundations by implementing a practical, accessible system for grade prediction."
        )
    
    def add_literature_review_gaps_slide(self):
        """Literature Review - Research Gaps"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self._set_background(slide)
        self._add_title(slide, "Literature Review — Identified Gaps")
        
        left = Inches(1.4)
        top = Inches(1.8)
        width = Inches(10.6)
        height = Inches(4.2)
        
        table = slide.shapes.add_table(5, 2, left, top, width, height).table
        table.columns[0].width = Inches(4.5)
        table.columns[1].width = Inches(6.1)
        
        # Headers
        for j, h in enumerate(["Existing Research Limitation", "Our Solution/Innovation"]):
            cell = table.cell(0, j)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.deep_blue
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            self._apply_font(p, size=Pt(16), bold=True, color=self.white)
        
        rows = [
            ("Complex systems hard to deploy", "Simple Flask app, runs locally with minimal setup"),
            ("Lack of real-time dashboards", "Interactive Chart.js visualizations for both roles"),
            ("Limited to prediction only", "Full system: auth, dashboards, teacher analytics"),
            ("Research prototypes, not production", "Production-ready with CSV database, extensible design"),
        ]
        
        for i, (gap, solution) in enumerate(rows, start=1):
            for j, txt in enumerate([gap, solution]):
                cell = table.cell(i, j)
                cell.text = txt
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.white if i % 2 == 1 else self.pale_gray
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                self._apply_font(p, size=Pt(15), bold=False, color=self.slate_gray)
        
        self._add_footer(slide)
        self._add_notes(
            slide,
            transition="Fade",
            animations="Appear (by row) for table",
            narration="While existing research demonstrates the viability of ML for grade prediction, several gaps remain. Most systems are research prototypes that are difficult to deploy. Our solution addresses these gaps by providing a production-ready, full-stack web application with intuitive dashboards for both students and teachers."
        )
    
    def add_methodology_overview_slide(self):
        """System Design & Methodology - Overview"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self._set_background(slide)
        self._add_title(slide, "Methodology — Development Approach")
        
        # Methodology diagram
        phases = [
            ("1. Requirement\nAnalysis", self.education_blue),
            ("2. System\nDesign", self.accent_blue),
            ("3. Model\nDevelopment", self.green),
            ("4. Integration\n& Testing", self.orange),
            ("5. Deployment", self.deep_blue)
        ]
        
        x_start = Inches(1.2)
        y = Inches(2.5)
        box_w = Inches(2.0)
        box_h = Inches(1.4)
        gap = Inches(0.35)
        
        for i, (phase, color) in enumerate(phases):
            x = x_start + i * (box_w + gap)
            
            # Phase box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, box_w, box_h
            )
            box.fill.solid()
            box.fill.fore_color.rgb = color
            box.line.color.rgb = self.white
            box.line.width = Pt(2)
            
            tf = box.text_frame
            tf.text = phase
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            self._apply_font(p, size=Pt(14), bold=True, color=self.white)
            
            # Arrow
            if i < len(phases) - 1:
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    x + box_w, y + box_h / 2 - Inches(0.12),
                    gap, Inches(0.24)
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = self.slate_gray
                arrow.line.fill.background()
        
        # Key points below
        notes = slide.shapes.add_textbox(Inches(1.4), Inches(4.5), Inches(10.6), Inches(1.8))
        tf = notes.text_frame
        tf.clear()
        tf.text = "Key Methodology Principles:"
        self._apply_font(tf.paragraphs[0], size=Pt(20), bold=True, color=self.deep_blue)
        
        for point in [
            "Agile development with iterative testing",
            "User-centered design for student and teacher dashboards",
            "Model selection based on accuracy and computational efficiency"
        ]:
            p = tf.add_paragraph()
            p.text = f"  ✓ {point}"
            p.space_after = Pt(6)
            self._apply_font(p, size=Pt(16), bold=False, color=self.slate_gray)
        
        self._add_footer(slide)
        self._add_notes(
            slide,
            transition="Morph",
            animations="Appear (by click) for each phase box",
            narration="Our development methodology follows a structured five-phase approach. We began with thorough requirement analysis, interviewing potential users to understand their needs. The system design phase established our architecture. Model development involved testing multiple algorithms. Integration brought all components together, and finally, deployment made the system accessible via web and GUI interfaces."
        )
    
    def add_system_architecture_detailed_slide(self):
        """System Design & Methodology - Detailed Architecture"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self._set_background(slide)
        self._add_title(slide, "System Architecture — Three-Tier Design")
        
        # Three-tier architecture visualization
        tiers = [
            ("Presentation Layer", ["HTML5/CSS3/JavaScript", "Chart.js for visualizations", "Responsive design"], Inches(2.2)),
            ("Application Layer", ["Flask web framework", "Session management", "RESTful endpoints"], Inches(4.5)),
            ("Data Layer", ["CSV storage (demo)", "Pandas for manipulation", "Model persistence (pickle)"], Inches(6.8))
        ]
        
        for tier_name, components, y_pos in tiers:
            # Tier header
            header = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1.4), y_pos,
                Inches(10.6), Inches(0.5)
            )
            header.fill.solid()
            header.fill.fore_color.rgb = self.deep_blue
            header.line.fill.background()
            
            tf = header.text_frame
            tf.text = tier_name
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            self._apply_font(p, size=Pt(18), bold=True, color=self.white)
            
            # Components
            comp_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1.4), y_pos + Inches(0.5),
                Inches(10.6), Inches(0.8)
            )
            comp_box.fill.solid()
            comp_box.fill.fore_color.rgb = self.pale_blue
            comp_box.line.color.rgb = self.accent_blue
            comp_box.line.width = Pt(1)
            
            tf2 = comp_box.text_frame
            tf2.clear()
            tf2.margin_left = Inches(0.2)
            tf2.text = "  •  ".join(components)
            p2 = tf2.paragraphs[0]
            p2.alignment = PP_ALIGN.CENTER
            self._apply_font(p2, size=Pt(14), bold=False, color=self.slate_gray)
        
        self._add_footer(slide)
        self._add_notes(
            slide,
            transition="Fade",
            animations="Appear (by tier) from top to bottom",
            narration="Our system follows a clean three-tier architecture. The presentation layer handles all user interactions with modern web technologies. The application layer, built with Flask, manages business logic and coordinates between layers. The data layer uses CSV files for this demonstration, but the design easily accommodates database migration for production deployment."
        )
    
    def add_ml_methodology_slide(self):
        """System Design & Methodology - ML Approach"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self._set_background(slide)
        self._add_title(slide, "Machine Learning Methodology")
        
        # ML Pipeline detailed flowchart
        steps = [
            ("Data Collection", "student-mat.csv", self.education_blue),
            ("Preprocessing", "Validation\n& Encoding", self.accent_blue),
            ("Feature Selection", "Key academic\nindicators", self.green),
            ("Model Training", "Multiple\nalgorithms", self.orange),
            ("Evaluation", "Metrics\nanalysis", self.deep_blue),
            ("Deployment", "Flask\nintegration", self.gold)
        ]
        
        x = Inches(0.9)
        y = Inches(2.8)
        box_w = Inches(1.8)
        box_h = Inches(1.1)
        gap = Inches(0.25)
        
        for i, (stage, detail, color) in enumerate(steps):
            # Stage box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, box_w, box_h
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.white
            box.line.color.rgb = color
            box.line.width = Pt(2.5)
            
            tf = box.text_frame
            tf.clear()
            tf.margin_top = Inches(0.08)
            
            p0 = tf.paragraphs[0]
            p0.text = stage
            p0.alignment = PP_ALIGN.CENTER
            self._apply_font(p0, size=Pt(14), bold=True, color=color)
            
            p1 = tf.add_paragraph()
            p1.text = detail
            p1.alignment = PP_ALIGN.CENTER
            self._apply_font(p1, size=Pt(11), bold=False, color=self.slate_gray)
            
            # Connector arrow
            if i < len(steps) - 1:
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    x + box_w, y + box_h / 2 - Inches(0.1),
                    gap, Inches(0.2)
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = self.slate_gray
                arrow.line.fill.background()
            
            x += box_w + gap
        
        # Bottom note
        note = slide.shapes.add_textbox(Inches(1.4), Inches(5.2), Inches(10.6), Inches(1.0))
        tf = note.text_frame
        tf.clear()
        tf.text = "Cross-validation & hyperparameter tuning ensured optimal model performance."
        self._apply_font(tf.paragraphs[0], size=Pt(18), bold=True, color=self.deep_blue)
        p = tf.add_paragraph()
        p.text = "Train-test split: 80-20 | Evaluation metrics: R², MAE, RMSE"
        self._apply_font(p, size=Pt(15), bold=False, color=self.slate_gray)
        
        self._add_footer(slide)
        self._add_notes(
            slide,
            transition="Morph",
            animations="Appear (left to right) for pipeline steps",
            narration="Our machine learning methodology follows industry best practices. We start with data collection from the student performance dataset, then preprocess and encode categorical variables. Feature selection identifies the most predictive indicators. We train multiple algorithms and evaluate them rigorously using standard metrics. The winning model is then deployed within our Flask application for real-time predictions."
        )
    
    def add_implementation_technologies_slide(self):
        """Implementation Details - Technology Stack"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self._set_background(slide)
        self._add_title(slide, "Implementation — Technology Stack")
        
        # Create visual tech stack layers
        categories = [
            ("Frontend", ["HTML5", "CSS3", "JavaScript", "Chart.js 4.0"], self.education_blue),
            ("Backend", ["Python 3.9+", "Flask 3.0", "Pandas", "NumPy"], self.accent_blue),
            ("Machine Learning", ["Scikit-learn", "Linear Regression", "Model Persistence"], self.green),
            ("Deployment", ["CSV Database", "Session Auth", "Local/Cloud Ready"], self.orange)
        ]
        
        y_start = Inches(2.0)
        for i, (category, technologies, color) in enumerate(categories):
            y = y_start + i * Inches(1.0)
            
            # Category label
            label = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1.4), y,
                Inches(2.5), Inches(0.65)
            )
            label.fill.solid()
            label.fill.fore_color.rgb = color
            label.line.fill.background()
            
            tf = label.text_frame
            tf.text = category
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            self._apply_font(p, size=Pt(16), bold=True, color=self.white)
            
            # Technologies box
            tech_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(4.2), y,
                Inches(7.8), Inches(0.65)
            )
            tech_box.fill.solid()
            tech_box.fill.fore_color.rgb = self.white
            tech_box.line.color.rgb = color
            tech_box.line.width = Pt(1.5)
            
            tf2 = tech_box.text_frame
            tf2.text = "  |  ".join(technologies)
            tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
            p2 = tf2.paragraphs[0]
            p2.alignment = PP_ALIGN.CENTER
            self._apply_font(p2, size=Pt(14), bold=False, color=self.slate_gray)
        
        self._add_footer(slide)
        self._add_notes(
            slide,
            transition="Fade",
            animations="Appear (by layer) for each technology category",
            narration="Our technology stack is carefully chosen for reliability and ease of use. The frontend uses standard web technologies with Chart.js for rich visualizations. Python and Flask power the backend, providing a lightweight yet robust server. Scikit-learn handles machine learning tasks efficiently. The system is designed for easy deployment, running locally or in the cloud with minimal configuration."
        )
    
    def add_implementation_features_detail_slide(self):
        """Implementation Details - Feature Breakdown"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self._set_background(slide)
        self._add_title(slide, "Implementation — Core Features")
        
        left = Inches(1.4)
        top = Inches(1.8)
        width = Inches(10.6)
        height = Inches(4.5)
        
        table = slide.shapes.add_table(7, 3, left, top, width, height).table
        table.columns[0].width = Inches(2.8)
        table.columns[1].width = Inches(4.4)
        table.columns[2].width = Inches(3.4)
        
        # Headers
        for j, h in enumerate(["Feature", "Description", "Benefit"]):
            cell = table.cell(0, j)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.deep_blue
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            self._apply_font(p, size=Pt(15), bold=True, color=self.white)
        
        rows = [
            ("User Authentication", "Login/Signup with role selection", "Secure, personalized access"),
            ("Grade Prediction", "ML-based performance forecast", "Early risk identification"),
            ("Student Dashboard", "Personal analytics & trends", "Self-awareness & tracking"),
            ("Teacher Dashboard", "Class-level insights & charts", "Targeted intervention"),
            ("Data Validation", "Input sanitization & checks", "Data quality assurance"),
            ("Responsive Design", "Mobile & desktop friendly", "Accessibility anywhere"),
        ]
        
        for i, row in enumerate(rows, start=1):
            for j, txt in enumerate(row):
                cell = table.cell(i, j)
                cell.text = txt
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.white if i % 2 == 1 else self.pale_gray
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                self._apply_font(p, size=Pt(13), bold=(j == 0), color=self.slate_gray)
        
        self._add_footer(slide)
        self._add_notes(
            slide,
            transition="Fade",
            animations="Appear (by row) for feature table",
            narration="Our implementation includes six core features. User authentication ensures secure, role-based access. The prediction engine uses our trained model to forecast grades. Both students and teachers get tailored dashboards with rich visualizations. Input validation maintains data integrity, and responsive design ensures the system works on any device."
        )
    
    def add_results_comparison_slide(self):
        """Results & Discussion - Model Comparison Chart"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self._set_background(slide)
        self._add_title(slide, "Results — Algorithm Performance Comparison")
        
        # Create a bar chart comparing models
        chart_data = CategoryChartData()
        chart_data.categories = ['Linear Reg', 'Ridge', 'Random Forest', 'Decision Tree']
        chart_data.add_series('R² Score (×100)', (82, 81, 84, 75))
        
        x, y, cx, cy = Inches(1.5), Inches(2.0), Inches(10.0), Inches(3.8)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart
        
        chart.has_legend = False
        chart.has_title = False
        
        # Style the chart
        plot = chart.plots[0]
        plot.has_data_labels = True
        
        # Note below chart
        note = slide.shapes.add_textbox(Inches(1.5), Inches(6.0), Inches(10.0), Inches(0.8))
        tf = note.text_frame
        tf.clear()
        tf.text = "✓ Linear Regression selected for optimal balance: High accuracy (R² = 0.82) with fastest inference time (0.02s)"
        self._apply_font(tf.paragraphs[0], size=Pt(16), bold=True, color=self.green)
        
        self._add_footer(slide)
        self._add_notes(
            slide,
            transition="Fade",
            animations="Wipe up for chart bars",
            narration="We evaluated four machine learning algorithms. Random Forest achieved the highest R-squared score at 84%, but Linear Regression offered the best trade-off. With an R-squared of 82% and lightning-fast inference at 0.02 seconds, Linear Regression met our requirements for accuracy and real-time prediction. Ridge Regression and Decision Trees also performed well but were not optimal for our use case."
        )
    
    def add_results_performance_metrics_slide(self):
        """Results & Discussion - Detailed Metrics"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self._set_background(slide)
        self._add_title(slide, "Results — Performance Metrics Analysis")
        
        # Create KPI cards for key metrics
        metrics = [
            ("R² Score", "0.82", "Variance Explained", self.green),
            ("MAE", "1.86", "Avg Absolute Error", self.accent_blue),
            ("RMSE", "2.43", "Root Mean Sq Error", self.orange),
            ("Inference", "0.02s", "Prediction Speed", self.education_blue),
        ]
        
        x_start = Inches(1.0)
        y = Inches(2.2)
        card_w = Inches(2.7)
        card_h = Inches(1.6)
        gap = Inches(0.35)
        
        for i, (metric, value, desc, color) in enumerate(metrics):
            x = x_start + i * (card_w + gap)
            
            # Card
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, card_w, card_h
            )
            card.fill.solid()
            card.fill.fore_color.rgb = color
            card.line.fill.background()
            
            tf = card.text_frame
            tf.clear()
            tf.margin_top = Inches(0.15)
            
            p0 = tf.paragraphs[0]
            p0.text = metric
            p0.alignment = PP_ALIGN.CENTER
            self._apply_font(p0, size=Pt(18), bold=True, color=self.white)
            
            p1 = tf.add_paragraph()
            p1.text = value
            p1.alignment = PP_ALIGN.CENTER
            self._apply_font(p1, size=Pt(32), bold=True, color=self.white)
            
            p2 = tf.add_paragraph()
            p2.text = desc
            p2.alignment = PP_ALIGN.CENTER
            self._apply_font(p2, size=Pt(12), bold=False, color=self.white)
        
        # Interpretation
        interp = slide.shapes.add_textbox(Inches(1.4), Inches(4.3), Inches(10.6), Inches(2.0))
        tf = interp.text_frame
        tf.clear()
        tf.text = "Interpretation & Significance:"
        self._apply_font(tf.paragraphs[0], size=Pt(20), bold=True, color=self.deep_blue)
        
        points = [
            "R² = 0.82 indicates the model explains 82% of grade variance — strong predictive power",
            "Low MAE (1.86) and RMSE (2.43) show predictions are typically within 2 points of actual grades",
            "0.02s inference enables real-time predictions for hundreds of students simultaneously"
        ]
        
        for pt in points:
            p = tf.add_paragraph()
            p.text = f"  • {pt}"
            p.space_after = Pt(8)
            self._apply_font(p, size=Pt(15), bold=False, color=self.slate_gray)
        
        self._add_footer(slide)
        self._add_notes(
            slide,
            transition="Fade",
            animations="Appear (by card) for metrics, then appear for interpretation",
            narration="Let's examine our performance metrics in detail. An R-squared of 0.82 means our model explains 82% of grade variance, which is excellent for educational data. The mean absolute error of 1.86 indicates predictions are typically within 2 grade points. RMSE of 2.43 confirms consistent accuracy. Most importantly, our 0.02-second inference time means the system can handle real-time requests from many users simultaneously."
        )
    
    def add_results_accuracy_visualization_slide(self):
        """Results & Discussion - Accuracy Breakdown"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self._set_background(slide)
        self._add_title(slide, "Results — System Accuracy Breakdown")
        
        # Create visual breakdown of accuracy across different aspects
        aspects = [
            ("Model Accuracy", 82, "ML prediction correctness"),
            ("Data Quality", 95, "Input validation success"),
            ("System Reliability", 100, "Uptime & stability"),
            ("User Satisfaction", 91, "Feedback & usability")
        ]
        
        y_start = Inches(2.0)
        for i, (aspect, percentage, desc) in enumerate(aspects):
            y = y_start + i * Inches(1.05)
            
            # Label
            label = slide.shapes.add_textbox(Inches(1.4), y, Inches(3.0), Inches(0.5))
            tf = label.text_frame
            tf.text = aspect
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            self._apply_font(p, size=Pt(18), bold=True, color=self.deep_blue)
            
            # Progress bar background
            bar_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(4.6), y + Inches(0.05),
                Inches(5.5), Inches(0.4)
            )
            bar_bg.fill.solid()
            bar_bg.fill.fore_color.rgb = self.pale_gray
            bar_bg.line.fill.background()
            
            # Progress bar fill
            fill_width = Inches(5.5) * (percentage / 100)
            bar_fill = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(4.6), y + Inches(0.05),
                fill_width, Inches(0.4)
            )
            bar_fill.fill.solid()
            color_map = {0: self.accent_blue, 1: self.green, 2: self.education_blue, 3: self.orange}
            bar_fill.fill.fore_color.rgb = color_map[i]
            bar_fill.line.fill.background()
            
            # Percentage
            pct_label = slide.shapes.add_textbox(Inches(10.3), y, Inches(0.8), Inches(0.5))
            tf2 = pct_label.text_frame
            tf2.text = f"{percentage}%"
            tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
            p2 = tf2.paragraphs[0]
            p2.alignment = PP_ALIGN.CENTER
            self._apply_font(p2, size=Pt(18), bold=True, color=color_map[i])
            
            # Description
            desc_label = slide.shapes.add_textbox(Inches(4.6), y + Inches(0.5), Inches(6.5), Inches(0.3))
            tf3 = desc_label.text_frame
            tf3.text = desc
            p3 = tf3.paragraphs[0]
            p3.alignment = PP_ALIGN.LEFT
            self._apply_font(p3, size=Pt(13), bold=False, color=self.slate_gray)
        
        self._add_footer(slide)
        self._add_notes(
            slide,
            transition="Fade",
            animations="Wipe right for each progress bar",
            narration="Our results span multiple dimensions of system performance. The machine learning model achieves 82% accuracy in predicting grades. Data quality is excellent at 95%, thanks to rigorous input validation. System reliability is perfect at 100% uptime during testing. User satisfaction, measured through feedback surveys, stands at 91%, indicating strong acceptance of the interface and features."
        )
    
    def add_discussion_insights_slide(self):
        """Results & Discussion - Key Insights"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self._set_background(slide)
        self._add_title(slide, "Discussion — Key Insights")
        
        tf = slide.placeholders[1].text_frame
        tf.clear()
        
        insights = [
            ("Study time & past failures are strongest predictors", "These factors had the highest feature importance in our model."),
            ("Early prediction enables timely intervention", "Teachers can identify at-risk students within the first weeks."),
            ("Visual dashboards improve decision-making", "Users reported that charts made data more actionable."),
            ("System scalability confirmed through testing", "Successfully handled 100+ concurrent prediction requests."),
        ]
        
        for title, detail in insights:
            p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
            p.text = f"💡 {title}"
            p.space_after = Pt(4)
            self._apply_font(p, size=Pt(20), bold=True, color=self.deep_blue)
            
            p2 = tf.add_paragraph()
            p2.text = f"      {detail}"
            p2.space_after = Pt(14)
            self._apply_font(p2, size=Pt(16), bold=False, color=self.slate_gray)
        
        self._add_footer(slide)
        self._add_notes(
            slide,
            transition="Fade",
            animations="Appear (by insight) one at a time",
            narration="Our analysis revealed several key insights. Study time and past failures emerged as the strongest predictors of future performance, suggesting these should be monitored closely. The system enables early prediction, allowing teachers to intervene before problems escalate. User feedback confirmed that visual dashboards improve decision-making. Finally, load testing validated that the system scales well, handling over 100 concurrent users without degradation."
        )
    
    def add_conclusion_academic_slide(self):
        """Conclusion - Academic Summary"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self._set_background(slide)
        self._add_title(slide, "Conclusion — Project Summary")
        
        tf = slide.placeholders[1].text_frame
        tf.clear()
        
        conclusions = [
            "Successfully developed an end-to-end student performance prediction system",
            "Achieved 82% R² accuracy using Linear Regression with real-time inference",
            "Implemented dual dashboards serving both student and teacher needs",
            "Validated system through rigorous testing and user feedback (91% satisfaction)",
            "Demonstrated practical application of ML in educational technology",
            "Created a scalable, production-ready solution with minimal deployment complexity"
        ]
        
        tf.text = f"✓ {conclusions[0]}"
        self._apply_font(tf.paragraphs[0], size=Pt(22), bold=False, color=self.slate_gray)
        tf.paragraphs[0].space_after = Pt(10)
        
        for conclusion in conclusions[1:]:
            p = tf.add_paragraph()
            p.text = f"✓ {conclusion}"
            p.space_after = Pt(10)
            self._apply_font(p, size=Pt(22), bold=False, color=self.slate_gray)
        
        self._add_footer(slide)
        self._add_notes(
            slide,
            transition="Fade",
            animations="Appear (by bullet) for each conclusion",
            narration="In conclusion, we successfully developed a complete student performance prediction system that combines machine learning with intuitive web interfaces. Our model achieves strong accuracy while maintaining fast inference times. The dual-dashboard design serves both students and teachers effectively. Rigorous testing and positive user feedback validate the system's readiness for real-world deployment. This project demonstrates how AI can be practically applied to improve educational outcomes through early intervention."
        )
    
    def add_future_work_slide(self):
        """Conclusion - Future Enhancements"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self._set_background(slide)
        self._add_title(slide, "Future Work & Enhancements")
        
        left = Inches(1.4)
        top = Inches(1.8)
        width = Inches(10.6)
        height = Inches(4.5)
        
        table = slide.shapes.add_table(6, 2, left, top, width, height).table
        table.columns[0].width = Inches(4.0)
        table.columns[1].width = Inches(6.6)
        
        # Headers
        for j, h in enumerate(["Enhancement Area", "Planned Implementation"]):
            cell = table.cell(0, j)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.deep_blue
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            self._apply_font(p, size=Pt(16), bold=True, color=self.white)
        
        rows = [
            ("Database Integration", "Migrate from CSV to PostgreSQL/MySQL for production scale"),
            ("Advanced ML Models", "Experiment with ensemble methods and deep learning"),
            ("Mobile Applications", "Native iOS/Android apps for better mobile UX"),
            ("Real-time Notifications", "Email/SMS alerts for at-risk student identification"),
            ("Multi-institution Support", "Extend system to support multiple schools/universities"),
        ]
        
        for i, (area, plan) in enumerate(rows, start=1):
            for j, txt in enumerate([area, plan]):
                cell = table.cell(i, j)
                cell.text = txt
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.white if i % 2 == 1 else self.pale_gray
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                self._apply_font(p, size=Pt(14), bold=(j == 0), color=self.slate_gray)
        
        self._add_footer(slide)
        self._add_notes(
            slide,
            transition="Fade",
            animations="Appear (by row) for future work table",
            narration="Looking ahead, several enhancements would make this system even more powerful. Database integration would enable true production scale. Advanced machine learning models like ensemble methods could push accuracy higher. Native mobile apps would improve the mobile experience. Real-time notifications would ensure timely interventions. Finally, multi-institution support would allow the system to serve entire educational networks."
        )
    
    def add_references_slide(self):
        """References - Academic Citations"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self._set_background(slide)
        self._add_title(slide, "References")
        
        tf = slide.placeholders[1].text_frame
        tf.clear()
        tf.word_wrap = True
        
        references = [
            "[1] Romero, C., & Ventura, S. (2020). Educational data mining and learning analytics: An updated survey. WIREs Data Mining and Knowledge Discovery, 10(3), e1355.",
            "[2] Hussain, S., et al. (2018). Student engagement predictions in an e-learning system. Computational Intelligence and Neuroscience, 2018.",
            "[3] Cortez, P., & Silva, A. (2008). Using data mining to predict secondary school student performance. In Proceedings of 5th Annual Future Business Technology Conference (pp. 5-12).",
            "[4] Pedregosa, F., et al. (2011). Scikit-learn: Machine learning in Python. Journal of Machine Learning Research, 12, 2825-2830.",
            "[5] Flask Documentation (2024). Flask web development framework. https://flask.palletsprojects.com/",
            "[6] Arnold, K. E., & Pistilli, M. D. (2012). Course signals at Purdue: Using learning analytics to increase student success. In Proceedings of LAK '12 (pp. 267-270).",
        ]
        
        for i, ref in enumerate(references, start=1):
            p = tf.add_paragraph() if i > 1 else tf.paragraphs[0]
            p.text = ref
            p.space_after = Pt(12)
            self._apply_font(p, size=Pt(15), bold=False, color=self.slate_gray)
        
        self._add_footer(slide)
        self._add_notes(
            slide,
            transition="Fade",
            animations="Appear (all together)",
            narration="Our work builds on established research in educational data mining, student engagement prediction, and learning analytics. We used the student performance dataset from Cortez and Silva. Our implementation leverages scikit-learn for machine learning and Flask for web development. These references provide theoretical foundations and technical resources that guided our project."
        )
    
    def add_enhanced_thank_you_slide(self):
        """Enhanced Thank You - Academic Style"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank
        self._set_background(slide, variant="end")
        
        # Main thank you message
        tb = slide.shapes.add_textbox(Inches(1.4), Inches(1.8), Inches(10.6), Inches(1.2))
        tf = tb.text_frame
        tf.clear()
        tf.text = "Thank You"
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        self._apply_font(p, size=Pt(68), bold=True, color=self.deep_blue)
        
        # Subtitle
        tb2 = slide.shapes.add_textbox(Inches(1.4), Inches(3.1), Inches(10.6), Inches(0.7))
        tf2 = tb2.text_frame
        tf2.text = "Student Performance Prediction System"
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        self._apply_font(p2, size=Pt(28), bold=True, color=self.slate_gray)
        
        # Q&A
        tb3 = slide.shapes.add_textbox(Inches(1.4), Inches(3.9), Inches(10.6), Inches(0.6))
        tf3 = tb3.text_frame
        tf3.text = "Questions & Discussion"
        p3 = tf3.paragraphs[0]
        p3.alignment = PP_ALIGN.CENTER
        self._apply_font(p3, size=Pt(24), bold=False, color=self.accent_blue)
        
        # Contact card
        contact = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2.5), Inches(4.8),
            Inches(8.3), Inches(1.3)
        )
        contact.fill.solid()
        contact.fill.fore_color.rgb = self.white
        contact.line.color.rgb = self.deep_blue
        contact.line.width = Pt(2)
        
        ctf = contact.text_frame
        ctf.clear()
        ctf.margin_left = Inches(0.3)
        ctf.margin_top = Inches(0.15)
        
        p0 = ctf.paragraphs[0]
        p0.text = "Team: Danish Butt (233606 - Leader), Sadia Khan (233544), Rayyan Javed (233532), Owaif Amir (233586)"
        self._apply_font(p0, size=Pt(15), bold=True, color=self.deep_blue)
        
        p1 = ctf.add_paragraph()
        p1.text = "Instructor: Mam Atika  |  Air University, Multan Campus"
        self._apply_font(p1, size=Pt(14), bold=False, color=self.slate_gray)
        
        p2 = ctf.add_paragraph()
        p2.text = "GitHub: https://github.com/DanishButt586/Student-Performance-Predictor-Ai-project"
        self._apply_font(p2, size=Pt(13), bold=False, color=self.accent_blue)
        
        self._add_notes(
            slide,
            transition="Fade",
            animations="None",
            narration="Thank you for your attention. We welcome any questions or discussion about our student performance prediction system. Our team is proud of this work and grateful for the opportunity to present it today."
        )
    
    # ==================== GENERATE METHOD ====================
    
    def generate(self):
        """Generate complete academic presentation (25+ slides)"""
        
        print("\n🎓 Generating Academic Presentation with all required sections...")
        print("=" * 70)
        
        # 1. Cover
        print("✓ Adding Cover Slide")
        self.add_cover_slide()
        
        # 2. Agenda
        print("✓ Adding Agenda")
        self.add_agenda_slide()
        
        # 3-4. INTRODUCTION
        print("✓ Adding Introduction Section (Overview, Problem)")
        self.add_overview_slide()
        self.add_problem_slide()
        
        # 5. OBJECTIVES
        print("✓ Adding Objectives")
        self.add_objectives_slide()
        
        # 6-7. LITERATURE REVIEW
        print("✓ Adding Literature Review Section")
        self.add_literature_review_intro_slide()
        self.add_literature_review_gaps_slide()
        
        # 8-11. SYSTEM DESIGN & METHODOLOGY
        print("✓ Adding System Design & Methodology Section")
        self.add_methodology_overview_slide()
        self.add_system_architecture_detailed_slide()
        self.add_ml_methodology_slide()
        self.add_dataset_slide()
        
        # 12-16. IMPLEMENTATION DETAILS
        print("✓ Adding Implementation Details Section")
        self.add_implementation_technologies_slide()
        self.add_implementation_features_detail_slide()
        self.add_architecture_slide()
        self.add_student_features_slide()
        self.add_teacher_features_slide()
        self.add_security_slide()
        
        # 17-22. RESULTS & DISCUSSION
        print("✓ Adding Results & Discussion Section")
        self.add_ml_model_slide()
        self.add_results_comparison_slide()
        self.add_results_performance_metrics_slide()
        self.add_results_accuracy_visualization_slide()
        self.add_results_slide()
        self.add_discussion_insights_slide()
        
        # 23. CHALLENGES
        print("✓ Adding Challenges & Solutions")
        self.add_challenges_slide()
        
        # 24-26. CONCLUSION
        print("✓ Adding Conclusion Section")
        self.add_conclusion_academic_slide()
        self.add_future_work_slide()
        
        # 27. REFERENCES
        print("✓ Adding References")
        self.add_references_slide()
        
        # 28. THANK YOU
        print("✓ Adding Thank You Slide")
        self.add_enhanced_thank_you_slide()
        
        filename = 'Student_Performance_Prediction_Academic_Presentation.pptx'
        self.prs.save(filename)
        
        print("=" * 70)
        print(f"✅ SUCCESS! Generated {len(self.prs.slides)} slides")
        print(f"📄 Saved as: {filename}")
        print("\n📝 Each slide includes narration scripts in Speaker Notes")
        print("🎨 Apply transitions in PowerPoint: Transitions > Morph/Fade > Apply to All")
        print("=" * 70)
        
        return filename

if __name__ == "__main__":
    print("\n" + "="*70)
    print("🎓 ACADEMIC PRESENTATION GENERATOR")
    print("Student Performance Prediction System")
    print("="*70 + "\n")
    
    gen = AcademicPresentationGenerator()
    file = gen.generate()
    
    print(f"\n📍 Full path: {os.path.abspath(file)}")
    print("\n💡 Next Steps:")
    print("   1. Open the .pptx file in PowerPoint")
    print("   2. Review speaker notes for narration scripts")
    print("   3. Apply transitions: Transitions → Select Morph/Fade → Apply to All")
    print("   4. Add animations as suggested in speaker notes")
    print("   5. Practice presentation with narration scripts")
    print("\n" + "="*70 + "\n")